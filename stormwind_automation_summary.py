from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Colors
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
ORANGE  = RGBColor(0xE8, 0x6B, 0x1F)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
GREEN   = RGBColor(0x1A, 0x8C, 0x4E)
BLUE    = RGBColor(0x1F, 0x6F, 0xC8)
RED     = RGBColor(0xC0, 0x39, 0x2B)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)

def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def para(tf, text, bold=False, size=11, color=DARK, italic=False, space_before=0, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def bullet_item(tf, bold_text, body_text=None, bullet_color=ORANGE, text_color=DARK, size=10, space_before=6):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    # Bullet character
    r0 = p.add_run()
    r0.text = "■  "
    r0.font.size = Pt(7)
    r0.font.color.rgb = bullet_color
    r0.font.bold = True
    # Bold label
    r1 = p.add_run()
    r1.text = bold_text
    r1.font.bold = True
    r1.font.size = Pt(size)
    r1.font.color.rgb = text_color
    if body_text:
        r2 = p.add_run()
        r2.text = "  " + body_text
        r2.font.bold = False
        r2.font.size = Pt(size)
        r2.font.color.rgb = text_color

def milestone_item(tf, date_text, desc, status, status_color, size=10, space_before=6):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    r0 = p.add_run()
    r0.text = "■  "
    r0.font.size = Pt(7)
    r0.font.color.rgb = ORANGE
    r0.font.bold = True
    if date_text:
        r1 = p.add_run()
        r1.text = date_text + "  "
        r1.font.bold = True
        r1.font.size = Pt(size)
        r1.font.color.rgb = DARK
    r2 = p.add_run()
    r2.text = desc + "  "
    r2.font.bold = False
    r2.font.size = Pt(size)
    r2.font.color.rgb = DARK
    r3 = p.add_run()
    r3.text = status
    r3.font.bold = True
    r3.font.size = Pt(size)
    r3.font.color.rgb = status_color

# ── Build slide ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# ── Navy header bar ──────────────────────────────────────────────────────────
header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.55))
header.fill.solid()
header.fill.fore_color.rgb = NAVY
header.line.fill.background()

# ── Title ────────────────────────────────────────────────────────────────────
title_box = add_text_box(slide, 0.3, 0.65, 12.7, 0.75)
tf = title_box.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r1 = p.add_run(); r1.text = "STORMWIND"
r1.font.bold = True; r1.font.size = Pt(32); r1.font.color.rgb = NAVY
r2 = p.add_run(); r2.text = " AUTOMATION SUMMARY"
r2.font.bold = False; r2.font.size = Pt(32); r2.font.color.rgb = DARK

# Sub-title date
sub_box = add_text_box(slide, 0.3, 1.28, 12.7, 0.3)
tf2 = sub_box.text_frame
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "MAY 2026"
r.font.size = Pt(11); r.font.color.rgb = MID_GRAY; r.font.bold = False

# ── Left column box (KEY OUTCOMES) ───────────────────────────────────────────
left_box_bg = slide.shapes.add_shape(1, Inches(0.3), Inches(1.7), Inches(6.1), Inches(4.6))
left_box_bg.fill.solid(); left_box_bg.fill.fore_color.rgb = RGBColor(0xF9,0xF9,0xF9)
left_box_bg.line.color.rgb = RGBColor(0xDD,0xDD,0xDD)

left = add_text_box(slide, 0.45, 1.8, 5.8, 4.4)
tf = left.text_frame; tf.word_wrap = True

para(tf, "KEY OUTCOMES", bold=True, size=13, color=ORANGE)

para(tf, "", size=5)
bullet_item(tf, "LMS App Suite (stormwind-automation)",
            "Playwright + POM. ~75 tests running nightly across admin & student flows.",
            size=9.5)

p = tf.add_paragraph(); p.space_before = Pt(3)
r = p.add_run(); r.text = "         Covers: login, course details, skills assessments, contact support, add users, my classroom & more."
r.font.size = Pt(9); r.font.color.rgb = MID_GRAY; r.font.italic = True

para(tf, "", size=4)
bullet_item(tf, "Foxbox Webpage Suite (foxboxwebpage-automation)",
            "Playwright + POM. 35 scenarios across the full foxbox.com site.",
            size=9.5)

p = tf.add_paragraph(); p.space_before = Pt(3)
r = p.add_run(); r.text = "         Covers: all service pages, AI assessment flow, link health, content cleanup & schema regression."
r.font.size = Pt(9); r.font.color.rgb = MID_GRAY; r.font.italic = True

para(tf, "", size=4)
bullet_item(tf, "Nightly CI via GitHub Actions",
            "Both suites run on schedule with Slack notifications and HTML reports.",
            size=9.5)

para(tf, "", size=4)
bullet_item(tf, "Reliability",
            "3 retries per test. Strict mode locators throughout. Fixtures handle pre-auth.",
            size=9.5)

# ── Right column: MILESTONES ─────────────────────────────────────────────────
right = add_text_box(slide, 6.65, 1.8, 6.35, 2.5)
tf = right.text_frame; tf.word_wrap = True

para(tf, "MILESTONES", bold=True, size=13, color=ORANGE)
para(tf, "", size=3)

milestone_item(tf, "Apr",  "Skills assessments strict mode fix", "FIXED", GREEN)
milestone_item(tf, "Apr",  "Contact support submit timing fix", "FIXED", GREEN)
milestone_item(tf, "May",  "FOX2-47 AI assessment questionnaire (SC30–SC33)", "COMPLETED", GREEN)
milestone_item(tf, "May",  "FOX2-40 Studio content cleanup validation (SC34)", "COMPLETED", GREEN)
milestone_item(tf, "May",  "FOX2-42 Schema bug fixes regression (SC35)", "COMPLETED", GREEN)
milestone_item(tf, "May",  "Approach & Send Ideas nightly fixes", "FIXED", GREEN)

# ── Right column: ATTENTION ITEMS ────────────────────────────────────────────
attn = add_text_box(slide, 6.65, 4.1, 6.35, 2.1)
tf = attn.text_frame; tf.word_wrap = True

para(tf, "ATTENTION ITEMS", bold=True, size=13, color=ORANGE)
para(tf, "", size=3)

bullet_item(tf, "FOX2-40 regression:", " 3 deleted service pages re-published (returning HTTP 200). Needs re-deletion in Sanity.", size=9.5)
para(tf, "", size=2)
bullet_item(tf, "Culture page broken link:", " /jobs returns HTTP 404. Link needs to be removed or updated in Sanity.", size=9.5)
para(tf, "", size=2)
bullet_item(tf, "LMS App defect:", " Remove-from-classroom flow broken — ADD TO CLASSROOM button never reappears. Ticket filed.", size=9.5)
para(tf, "", size=2)
bullet_item(tf, "Flaky (ongoing):", " Admin add-users Leaderboard link intermittently not found. Passes on retry.", size=9.5)

# ── Bottom info boxes ─────────────────────────────────────────────────────────
for x, label, value in [(0.3, "APP SUITE", "stormwind-automation"), (3.5, "WEBPAGE SUITE", "foxboxwebpage-automation"), (7.2, "TOTAL TESTS", "~110 tests")]:
    box_bg = slide.shapes.add_shape(1, Inches(x), Inches(6.55), Inches(2.9), Inches(0.72))
    box_bg.fill.solid(); box_bg.fill.fore_color.rgb = GRAY
    box_bg.line.color.rgb = RGBColor(0xCC,0xCC,0xCC)
    tb = add_text_box(slide, x + 0.08, 6.57, 2.75, 0.65)
    tf = tb.text_frame
    para(tf, label, bold=False, size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)
    para(tf, value, bold=True, size=11, color=NAVY, align=PP_ALIGN.CENTER)

# ── Footer branding ───────────────────────────────────────────────────────────
footer = add_text_box(slide, 10.8, 6.6, 2.3, 0.4)
tf = footer.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
r = p.add_run(); r.text = "Foxbox Digital"
r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = NAVY

out = "/Users/santiagodalmas/work/stormwind_playwright_project/stormwind_automation_summary.pptx"
prs.save(out)
print(f"Saved: {out}")
